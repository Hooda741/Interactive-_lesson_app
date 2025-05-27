import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional

class PPTXGenerator:
    """
    Service for generating PowerPoint presentations from extracted content
    """
    
    def __init__(self, templates_folder: str, output_folder: str):
        """
        Initialize the PPTX generator
        
        Args:
            templates_folder: Path to the folder containing presentation templates
            output_folder: Path to the folder where generated presentations will be saved
        """
        self.templates_folder = templates_folder
        self.output_folder = output_folder
        
        # Ensure folders exist
        if not os.path.exists(templates_folder):
            os.makedirs(templates_folder)
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
        
        # Default colors for different templates
        self.color_schemes = {
            "default": {
                "title": RGBColor(89, 49, 150),  # Purple
                "heading": RGBColor(0, 112, 192),  # Blue
                "subheading": RGBColor(0, 176, 80),  # Green
                "text": RGBColor(0, 0, 0),  # Black
                "background": RGBColor(255, 255, 255),  # White
                "accent": RGBColor(255, 192, 0)  # Yellow
            },
            "colorful": {
                "title": RGBColor(192, 0, 0),  # Red
                "heading": RGBColor(0, 112, 192),  # Blue
                "subheading": RGBColor(112, 48, 160),  # Purple
                "text": RGBColor(0, 0, 0),  # Black
                "background": RGBColor(255, 255, 255),  # White
                "accent": RGBColor(255, 192, 0)  # Yellow
            },
            "minimal": {
                "title": RGBColor(0, 0, 0),  # Black
                "heading": RGBColor(68, 68, 68),  # Dark Gray
                "subheading": RGBColor(102, 102, 102),  # Gray
                "text": RGBColor(0, 0, 0),  # Black
                "background": RGBColor(255, 255, 255),  # White
                "accent": RGBColor(0, 112, 192)  # Blue
            }
        }
    
    def generate_presentation(self, content: Dict[str, Any], template_name: str = "default", 
                             color_scheme: str = "default", output_filename: str = None) -> str:
        """
        Generate a PowerPoint presentation from extracted content
        
        Args:
            content: Dictionary containing extracted content
            template_name: Name of the template to use
            color_scheme: Name of the color scheme to use
            output_filename: Name of the output file (without extension)
            
        Returns:
            Path to the generated presentation
        """
        # Determine template path
        template_path = os.path.join(self.templates_folder, f"{template_name}.pptx")
        
        # If template doesn't exist, use a blank presentation
        if not os.path.exists(template_path):
            prs = Presentation()
        else:
            prs = Presentation(template_path)
        
        # Set slide dimensions for standard 4:3 presentation
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Get color scheme
        colors = self.color_schemes.get(color_scheme, self.color_schemes["default"])
        
        # Generate title slide
        self._add_title_slide(prs, content, colors)
        
        # Generate content slides
        self._add_content_slides(prs, content, colors)
        
        # Generate activity slides
        self._add_activity_slides(prs, content, colors)
        
        # Generate summary slide
        self._add_summary_slide(prs, content, colors)
        
        # Determine output path
        if output_filename is None:
            output_filename = f"presentation_{os.path.basename(content.get('title', 'untitled'))}"
        
        output_path = os.path.join(self.output_folder, f"{output_filename}.pptx")
        
        # Save the presentation
        prs.save(output_path)
        
        return output_path
    
    def _add_title_slide(self, prs: Presentation, content: Dict[str, Any], colors: Dict[str, RGBColor]):
        """
        Add a title slide to the presentation
        
        Args:
            prs: Presentation object
            content: Dictionary containing extracted content
            colors: Dictionary containing color scheme
        """
        # Get title slide layout
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = content.get("title", "Untitled Presentation")
        if "structure" in content and "title" in content["structure"] and content["structure"]["title"]:
            title = content["structure"]["title"]
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        
        # Set subtitle (if available)
        if hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
            subtitle = "Interactive Lesson"
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_para = subtitle_shape.text_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = colors["heading"]
    
    def _add_content_slides(self, prs: Presentation, content: Dict[str, Any], colors: Dict[str, RGBColor]):
        """
        Add content slides to the presentation
        
        Args:
            prs: Presentation object
            content: Dictionary containing extracted content
            colors: Dictionary containing color scheme
        """
        # Get content slide layout
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        
        # Add slides for headings
        if "structure" in content and "headings" in content["structure"]:
            for heading in content["structure"]["headings"]:
                slide = prs.slides.add_slide(slide_layout)
                
                # Set heading as slide title
                title_shape = slide.shapes.title
                title_shape.text = heading["text"]
                title_para = title_shape.text_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                title_para.font.size = Pt(40)
                title_para.font.bold = True
                title_para.font.color.rgb = colors["title"]
                
                # Find paragraphs related to this heading
                related_paragraphs = self._find_related_paragraphs(content, heading)
                
                # Add paragraphs to slide
                if related_paragraphs and hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
                    content_shape = slide.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.clear()  # Clear default text
                    
                    for paragraph in related_paragraphs:
                        p = text_frame.add_paragraph()
                        p.text = paragraph
                        p.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                        p.font.size = Pt(24)
                        p.font.color.rgb = colors["text"]
        
        # Add slides for bullet points
        if "structure" in content and "bullet_points" in content["structure"]:
            for bullet_list in content["structure"]["bullet_points"]:
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title_shape = slide.shapes.title
                title_shape.text = "النقاط الرئيسية"  # "Key Points" in Arabic
                title_para = title_shape.text_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                title_para.font.size = Pt(40)
                title_para.font.bold = True
                title_para.font.color.rgb = colors["title"]
                
                # Add bullet points to slide
                if hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
                    content_shape = slide.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.clear()  # Clear default text
                    
                    for bullet_point in bullet_list:
                        p = text_frame.add_paragraph()
                        p.text = bullet_point
                        p.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                        p.font.size = Pt(24)
                        p.level = 0
                        p.font.color.rgb = colors["text"]
        
        # Add slides for images
        if "images" in content and content["images"]:
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "الصور التوضيحية"  # "Illustrations" in Arabic
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = colors["title"]
            
            # Add images to slide
            # For simplicity, we'll just add the first image
            if content["images"] and "path" in content["images"][0]:
                image_path = content["images"][0]["path"]
                if os.path.exists(image_path):
                    left = Inches(2)
                    top = Inches(2)
                    width = Inches(6)
                    height = Inches(4)
                    
                    slide.shapes.add_picture(image_path, left, top, width, height)
    
    def _add_activity_slides(self, prs: Presentation, content: Dict[str, Any], colors: Dict[str, RGBColor]):
        """
        Add activity slides to the presentation
        
        Args:
            prs: Presentation object
            content: Dictionary containing extracted content
            colors: Dictionary containing color scheme
        """
        # Get activity slide layout
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        
        # Add a sample activity slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "نشاط تفاعلي"  # "Interactive Activity" in Arabic
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        
        # Add activity description
        if hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default text
            
            p = text_frame.add_paragraph()
            p.text = "أكمل النشاط التالي:"  # "Complete the following activity:" in Arabic
            p.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.color.rgb = colors["heading"]
            
            # Add activity steps
            activity_steps = [
                "اقرأ النص بعناية",  # "Read the text carefully" in Arabic
                "حدد المفاهيم الرئيسية",  # "Identify the main concepts" in Arabic
                "أجب عن الأسئلة",  # "Answer the questions" in Arabic
                "ناقش إجاباتك مع زملائك"  # "Discuss your answers with your classmates" in Arabic
            ]
            
            for step in activity_steps:
                p = text_frame.add_paragraph()
                p.text = step
                p.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                p.font.size = Pt(24)
                p.level = 1
                p.font.color.rgb = colors["text"]
    
    def _add_summary_slide(self, prs: Presentation, content: Dict[str, Any], colors: Dict[str, RGBColor]):
        """
        Add a summary slide to the presentation
        
        Args:
            prs: Presentation object
            content: Dictionary containing extracted content
            colors: Dictionary containing color scheme
        """
        # Get summary slide layout
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "ملخص الدرس"  # "Lesson Summary" in Arabic
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        
        # Add summary points
        if hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default text
            
            # Extract key points from content
            summary_points = self._generate_summary_points(content)
            
            for point in summary_points:
                p = text_frame.add_paragraph()
                p.text = point
                p.alignment = PP_ALIGN.RIGHT  # Right-aligned for Arabic
                p.font.size = Pt(24)
                p.level = 0
                p.font.color.rgb = colors["text"]
    
    def _find_related_paragraphs(self, content: Dict[str, Any], heading: Dict[str, Any]) -> List[str]:
        """
        Find paragraphs related to a heading
        
        Args:
            content: Dictionary containing extracted content
            heading: Dictionary containing heading information
            
        Returns:
            List of related paragraphs
        """
        related_paragraphs = []
        
        if "structure" not in content or "paragraphs" not in content["structure"]:
            return related_paragraphs
        
        # Simple heuristic: take the first 2 paragraphs after the heading
        # In a real application, this would be more sophisticated
        heading_text = heading["text"]
        heading_index = -1
        
        for i, h in enumerate(content["structure"]["headings"]):
            if h["text"] == heading_text:
                heading_index = i
                break
        
        if heading_index != -1 and heading_index < len(content["structure"]["headings"]) - 1:
            next_heading_index = heading_index + 1
            next_heading = content["structure"]["headings"][next_heading_index]
            
            # Find paragraphs between this heading and the next
            for paragraph in content["structure"]["paragraphs"]:
                # Simple heuristic: if paragraph contains the heading text, it's related
                if heading_text in paragraph and next_heading["text"] not in paragraph:
                    related_paragraphs.append(paragraph)
        else:
            # If this is the last heading, take all remaining paragraphs
            for paragraph in content["structure"]["paragraphs"]:
                if heading_text in paragraph:
                    related_paragraphs.append(paragraph)
        
        # Limit to first 2 paragraphs to avoid overcrowding the slide
        return related_paragraphs[:2]
    
    def _generate_summary_points(self, content: Dict[str, Any]) -> List[str]:
        """
        Generate summary points from content
        
        Args:
            content: Dictionary containing extracted content
            
        Returns:
            List of summary points
        """
        summary_points = []
        
        # Extract headings as summary points
        if "structure" in content and "headings" in content["structure"]:
            for heading in content["structure"]["headings"]:
                summary_points.append(heading["text"])
        
        # Extract first bullet point from each bullet list
        if "structure" in content and "bullet_points" in content["structure"]:
            for bullet_list in content["structure"]["bullet_points"]:
                if bullet_list:
                    summary_points.append(bullet_list[0])
        
        # Limit to 5 points to avoid overcrowding the slide
        return summary_points[:5]
